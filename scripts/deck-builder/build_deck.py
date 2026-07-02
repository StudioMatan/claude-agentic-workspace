#!/usr/bin/env python3
"""
build_deck.py
--------------
Generic, reusable slide-deck builder on top of python-pptx.

Provides a small set of primitives (slide / box / txt / stat / title) styled
with a navy/amber visual system, plus a set of composed slide-type functions
(title slide, stat-row slide, two-column before/after flow, numbered list,
table, closing slide). None of it is org-specific - every string, color,
and number comes from a data dict (or an equivalent JSON file) passed in
by the caller.

Usage as a library:

    from build_deck import DeckBuilder

    deck = DeckBuilder()
    deck.title_slide("My Project", "One-line subtitle", "footer text")
    deck.stat_slide("Title", "Subtitle", [("42", "label"), ...])
    deck.save("out.pptx")

Usage as a CLI (build from a JSON spec):

    python3 build_deck.py --spec deck.json --out deck.pptx

Run with no arguments to build the demonstration deck bundled in __main__
below (the Falcon -> Claude SOAR overview deck) - this doubles as the
worked example for how to drive the library from a data dict.
"""
from __future__ import annotations

import argparse
import json
import os
from typing import Iterable, Optional, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Visual design system - navy / amber, generic (no org branding)
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
INK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xF2, 0xA9, 0x00)
RISK = RGBColor(0xC0, 0x39, 0x2B)
GOOD = RGBColor(0x2E, 0x7D, 0x32)
CARD = RGBColor(0xFF, 0xFF, 0xFF)

HDR_FONT = "Georgia"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


class DeckBuilder:
    """Thin, generic wrapper around python-pptx with a consistent visual system.

    All content is passed in by the caller - this class has no knowledge of
    any specific organization, project, or dataset.
    """

    def __init__(self, width=SLIDE_W, height=SLIDE_H):
        self.prs = Presentation()
        self.prs.slide_width = width
        self.prs.slide_height = height
        self._blank = self.prs.slide_layouts[6]
        self.width = width
        self.height = height

    # -- primitives ---------------------------------------------------

    def slide(self, bg: RGBColor = LIGHT):
        """New blank slide filled with a solid background color."""
        s = self.prs.slides.add_slide(self._blank)
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.width, self.height)
        r.fill.solid()
        r.fill.fore_color.rgb = bg
        r.line.fill.background()
        r.shadow.inherit = False
        s.shapes._spTree.remove(r._element)
        s.shapes._spTree.insert(2, r._element)
        return s

    def box(self, s, l, t, w, h, fill: Optional[RGBColor] = None, radius=False):
        """Rectangle or rounded-rectangle shape, used as a card/divider/bar."""
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shp = s.shapes.add_shape(shape_type, l, t, w, h)
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def txt(self, s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
        """Multi-paragraph, multi-run text box.

        `runs` is a list of lines: [{"segs": [{"t": text, "f": font, "s": size,
        "b": bold, "i": italic, "c": color}, ...]}, ...]. Only "t" is required
        per segment; everything else has a sane default.
        """
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        first = True
        for line in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.line_spacing = sp
            for seg in line["segs"]:
                r = p.add_run()
                r.text = seg["t"]
                r.font.name = seg.get("f", BODY_FONT)
                r.font.size = Pt(seg["s"] if "s" in seg else 14)
                r.font.bold = seg.get("b", False)
                r.font.italic = seg.get("i", False)
                r.font.color.rgb = seg.get("c", INK)
        return tb

    def stat(self, s, l, t, w, num, label, numc=NAVY):
        """Single stat card: big number over a small muted label."""
        self.box(s, l, t, w, Inches(1.7), fill=CARD, radius=True)
        self.txt(
            s, l + Inches(0.1), t + Inches(0.18), w - Inches(0.2), Inches(0.95),
            [{"segs": [{"t": num, "f": HDR_FONT, "s": 36, "b": True, "c": numc}]}],
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
        )
        self.txt(
            s, l + Inches(0.12), t + Inches(1.12), w - Inches(0.24), Inches(0.5),
            [{"segs": [{"t": label, "s": 12.5, "c": MUTED}]}],
            PP_ALIGN.CENTER,
        )

    def slide_title(self, s, text, sub=None, dark=False):
        """Standard slide heading + optional subtitle, top-left."""
        self.txt(
            s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.9),
            [{"segs": [{"t": text, "f": HDR_FONT, "s": 30, "b": True, "c": WHITE if dark else NAVY}]}],
        )
        if sub:
            self.txt(
                s, Inches(0.72), Inches(1.22), Inches(12), Inches(0.5),
                [{"segs": [{"t": sub, "s": 14.5, "c": AMBER if dark else MUTED}]}],
            )

    # -- composed slide types ------------------------------------------

    def title_slide(self, title: str, subtitle: str = "", footer: str = ""):
        s = self.slide(NAVY)
        self.box(s, Inches(0.7), Inches(2.55), Inches(0.18), Inches(2.0), fill=AMBER)
        self.txt(
            s, Inches(1.05), Inches(2.5), Inches(11.2), Inches(1.6),
            [{"segs": [{"t": title, "f": HDR_FONT, "s": 46, "b": True, "c": WHITE}]}],
        )
        if subtitle:
            self.txt(
                s, Inches(1.07), Inches(3.95), Inches(11), Inches(0.9),
                [{"segs": [{"t": subtitle, "s": 19, "c": ICE}]}], sp=1.15,
            )
        if footer:
            self.txt(
                s, Inches(1.07), Inches(6.55), Inches(11), Inches(0.5),
                [{"segs": [{"t": footer, "s": 13, "c": AMBER}]}],
            )
        return s

    def stat_slide(self, title: str, subtitle: str, stats: Sequence[tuple], body: str = ""):
        """Title + a row of up to 4 stat cards + an optional body paragraph."""
        s = self.slide(LIGHT)
        self.slide_title(s, title, subtitle)
        n = len(stats)
        gap = Inches(0.22)
        total_w = Inches(12.0)
        w = Inches((total_w.inches - gap.inches * (n - 1)) / n) if n else total_w
        x0 = Inches(0.7)
        y = Inches(2.5)
        for i, (num, label) in enumerate(stats):
            self.stat(s, x0 + i * (w + gap), y, w, num, label)
        if body:
            self.txt(
                s, Inches(0.7), Inches(4.75), Inches(12), Inches(1.6),
                [{"segs": [{"t": body, "s": 15.5, "c": INK}]}], sp=1.2,
            )
        return s

    def flow_compare_slide(
        self, title: str, subtitle: str,
        left_label: str, left_steps: Sequence[str],
        right_label: str, right_steps: Sequence[str],
        callout: str = "",
    ):
        """Two-column step-by-step flow comparison (e.g. before/after).

        Renders each side as a numbered vertical chain of steps inside a
        card, side by side, with an optional callout strip beneath
        explaining what structurally moved between the two paths.
        """
        s = self.slide(LIGHT)
        self.slide_title(s, title, subtitle)

        col_w = Inches(5.85)
        gap = Inches(0.3)
        left_x = Inches(0.7)
        right_x = left_x + col_w + gap
        top_y = Inches(2.15)
        card_h = Inches(4.05)

        def render_column(x, label, steps, accent):
            self.box(s, x, top_y, col_w, card_h, fill=CARD, radius=True)
            self.box(s, x, top_y, col_w, Inches(0.55), fill=accent, radius=True)
            self.txt(
                s, x + Inches(0.25), top_y, col_w - Inches(0.5), Inches(0.55),
                [{"segs": [{"t": label, "f": HDR_FONT, "s": 16, "b": True, "c": WHITE}]}],
                anchor=MSO_ANCHOR.MIDDLE,
            )
            step_h = (card_h.inches - 0.75) / max(len(steps), 1)
            yy = top_y + Inches(0.72)
            for i, step in enumerate(steps, start=1):
                bullet_d = Inches(0.32)
                self.box(s, x + Inches(0.25), yy, bullet_d, bullet_d, fill=accent, radius=True)
                self.txt(
                    s, x + Inches(0.25), yy, bullet_d, bullet_d,
                    [{"segs": [{"t": str(i), "f": HDR_FONT, "s": 13, "b": True, "c": WHITE}]}],
                    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
                )
                self.txt(
                    s, x + Inches(0.72), yy - Inches(0.03), col_w - Inches(1.0), Inches(step_h),
                    [{"segs": [{"t": step, "s": 12.5, "c": INK}]}], sp=1.05,
                )
                yy += Inches(step_h)
            return s

        render_column(left_x, left_label, left_steps, MUTED)
        render_column(right_x, right_label, right_steps, NAVY)

        if callout:
            self.box(s, Inches(0.7), Inches(6.55), Inches(11.95), Inches(0.65), fill=NAVY, radius=True)
            self.txt(
                s, Inches(0.95), Inches(6.55), Inches(11.5), Inches(0.65),
                [{"segs": [{"t": callout, "s": 13, "b": True, "c": AMBER}]}],
                anchor=MSO_ANCHOR.MIDDLE, sp=1.05,
            )
        return s

    def pipeline_slide(self, title: str, subtitle: str, steps: Sequence[tuple], caption: str = ""):
        """Vertical numbered pipeline (architecture / process flow).

        `steps` is a list of (heading, description) tuples, rendered top to
        bottom as numbered rows with a connecting rail.
        """
        s = self.slide(LIGHT)
        self.slide_title(s, title, subtitle)
        y = Inches(2.15)
        row_h = Inches(0.92)
        for i, (head, desc) in enumerate(steps, start=1):
            circle = Inches(0.5)
            self.box(s, Inches(0.7), y, circle, circle, fill=NAVY, radius=True)
            self.txt(
                s, Inches(0.7), y, circle, circle,
                [{"segs": [{"t": str(i), "f": HDR_FONT, "s": 18, "b": True, "c": WHITE}]}],
                PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
            )
            self.txt(
                s, Inches(1.45), y - Inches(0.02), Inches(11.0), Inches(0.4),
                [{"segs": [{"t": head, "f": HDR_FONT, "s": 16.5, "b": True, "c": NAVY}]}],
            )
            self.txt(
                s, Inches(1.45), y + Inches(0.38), Inches(11.0), Inches(0.5),
                [{"segs": [{"t": desc, "s": 13, "c": INK}]}], sp=1.05,
            )
            y += row_h
        if caption:
            self.txt(
                s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
                [{"segs": [{"t": caption, "s": 12.5, "i": True, "c": MUTED}]}],
            )
        return s

    def numbered_list_slide(self, title: str, subtitle: str, items: Sequence[tuple], dark=False):
        """Numbered list of (heading, body) items - e.g. design principles."""
        s = self.slide(NAVY if dark else LIGHT)
        self.slide_title(s, title, subtitle, dark=dark)
        y = Inches(2.15)
        row_h = Inches(1.55)
        head_c = AMBER if dark else NAVY
        body_c = ICE if dark else INK
        for i, (head, body) in enumerate(items, start=1):
            circle = Inches(0.55)
            self.box(s, Inches(0.7), y, circle, circle, fill=AMBER, radius=True)
            self.txt(
                s, Inches(0.7), y, circle, circle,
                [{"segs": [{"t": str(i), "f": HDR_FONT, "s": 20, "b": True, "c": NAVY}]}],
                PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
            )
            self.txt(
                s, Inches(1.5), y - Inches(0.02), Inches(11.0), Inches(0.45),
                [{"segs": [{"t": head, "f": HDR_FONT, "s": 17, "b": True, "c": head_c}]}],
            )
            self.txt(
                s, Inches(1.5), y + Inches(0.42), Inches(10.9), Inches(0.95),
                [{"segs": [{"t": body, "s": 13.5, "c": body_c}]}], sp=1.1,
            )
            y += row_h
        return s

    def table_slide(self, title: str, subtitle: str, headers: Sequence[str], rows: Sequence[Sequence[str]]):
        """Simple styled table slide (e.g. phase status)."""
        s = self.slide(LIGHT)
        self.slide_title(s, title, subtitle)
        n_cols = len(headers)
        table_l, table_t = Inches(0.7), Inches(2.3)
        table_w, table_h = Inches(11.95), Inches(0.6) + Inches(0.9) * len(rows)
        gshape = s.shapes.add_table(len(rows) + 1, n_cols, table_l, table_t, table_w, table_h)
        table = gshape.table
        # header row
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.text_frame.text = h
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(13)
                p.font.color.rgb = WHITE
                p.font.name = BODY_FONT
        # body rows
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD if r % 2 else LIGHT
                cell.text_frame.word_wrap = True
                cell.text_frame.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = INK
                    p.font.name = BODY_FONT
        return s

    def closing_slide(self, title: str, subtitle: str, links: Sequence[str]):
        s = self.slide(NAVY)
        self.box(s, Inches(0.7), Inches(2.3), Inches(0.18), Inches(1.6), fill=AMBER)
        self.txt(
            s, Inches(1.05), Inches(2.25), Inches(11), Inches(1.1),
            [{"segs": [{"t": title, "f": HDR_FONT, "s": 38, "b": True, "c": WHITE}]}],
        )
        self.txt(
            s, Inches(1.07), Inches(3.35), Inches(11), Inches(0.6),
            [{"segs": [{"t": subtitle, "s": 16, "c": ICE}]}],
        )
        y = Inches(4.4)
        for link in links:
            self.txt(
                s, Inches(1.07), y, Inches(11), Inches(0.45),
                [{"segs": [{"t": link, "s": 15, "b": True, "c": AMBER}]}],
            )
            y += Inches(0.55)
        return s

    def body_slide(self, title: str, subtitle: str, body: str):
        """Plain title + paragraph slide."""
        s = self.slide(LIGHT)
        self.slide_title(s, title, subtitle)
        self.txt(
            s, Inches(0.7), Inches(2.3), Inches(12), Inches(4.0),
            [{"segs": [{"t": body, "s": 15.5, "c": INK}]}], sp=1.25,
        )
        return s

    # -- output ---------------------------------------------------------

    def save(self, path: str):
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        self.prs.save(path)
        return path

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides.__iter__.__self__._sldIdLst)


# ---------------------------------------------------------------------------
# JSON-spec driven build (generic CLI entrypoint)
# ---------------------------------------------------------------------------

SLIDE_BUILDERS = {
    "title": lambda d, sl: d.title_slide(sl.get("title", ""), sl.get("subtitle", ""), sl.get("footer", "")),
    "stats": lambda d, sl: d.stat_slide(sl.get("title", ""), sl.get("subtitle", ""), [tuple(x) for x in sl.get("stats", [])], sl.get("body", "")),
    "flow_compare": lambda d, sl: d.flow_compare_slide(
        sl.get("title", ""), sl.get("subtitle", ""),
        sl.get("left_label", ""), sl.get("left_steps", []),
        sl.get("right_label", ""), sl.get("right_steps", []),
        sl.get("callout", ""),
    ),
    "pipeline": lambda d, sl: d.pipeline_slide(sl.get("title", ""), sl.get("subtitle", ""), [tuple(x) for x in sl.get("steps", [])], sl.get("caption", "")),
    "numbered_list": lambda d, sl: d.numbered_list_slide(sl.get("title", ""), sl.get("subtitle", ""), [tuple(x) for x in sl.get("items", [])], sl.get("dark", False)),
    "table": lambda d, sl: d.table_slide(sl.get("title", ""), sl.get("subtitle", ""), sl.get("headers", []), sl.get("rows", [])),
    "body": lambda d, sl: d.body_slide(sl.get("title", ""), sl.get("subtitle", ""), sl.get("body", "")),
    "closing": lambda d, sl: d.closing_slide(sl.get("title", ""), sl.get("subtitle", ""), sl.get("links", [])),
}


def build_from_spec(spec: dict) -> DeckBuilder:
    """Build a full deck from a generic JSON-serializable spec:

        {"slides": [{"type": "title", "title": "...", ...}, ...]}

    Each slide dict's "type" selects the composed builder function; the rest
    of the dict is passed through as kwargs-equivalent fields. No slide type
    here knows anything about a specific organization or dataset - all
    content is data supplied by the caller.
    """
    deck = DeckBuilder()
    for sl in spec.get("slides", []):
        kind = sl.get("type")
        fn = SLIDE_BUILDERS.get(kind)
        if fn is None:
            raise ValueError(f"Unknown slide type: {kind!r}")
        fn(deck, sl)
    return deck


def _cli():
    ap = argparse.ArgumentParser(description="Build a slide deck from a generic JSON spec.")
    ap.add_argument("--spec", help="Path to a JSON spec file ({'slides': [...]})")
    ap.add_argument("--out", default=None, help="Output .pptx path (default: deck.pptx for --spec builds, or the bundled demo's default path)")
    args = ap.parse_args()

    if args.spec:
        with open(args.spec) as f:
            spec = json.load(f)
        deck = build_from_spec(spec)
        path = deck.save(args.out or "deck.pptx")
        print(f"Saved: {path} ({deck.slide_count} slides)")
    else:
        _build_demo_deck(args.out)


# ---------------------------------------------------------------------------
# Demonstration build: the Falcon -> Claude SOAR overview deck
#
# This is the worked example showing how a real deck is driven from a data
# dict using only the generic primitives above. Content mirrors
# ../../flows/falcon-claude-soar/deck-content.md - edit that file first if
# you're changing what the deck says, then keep this dict in sync.
# ---------------------------------------------------------------------------

def _build_demo_deck(out_path: Optional[str] = None) -> str:
    deck = DeckBuilder()

    # 1 - title
    deck.title_slide(
        "Falcon → Claude: Automated Detection Triage",
        "A structured verdict on every EDR detection, in under a minute - "
        "human reviews, doesn't assemble",
        "CrowdStrike Falcon + Claude (Anthropic) · illustrative deployment",
    )

    # 2 - the problem
    deck.stat_slide(
        "Manual triage doesn't scale",
        "Not a volume problem - a speed, documentation, and consistency problem",
        [
            ("~1.5 days", "median wall-clock time, detection to close"),
            ("<10%", "of detections get a written report"),
            ("~30 min", "hands-on analyst time per detection reviewed"),
        ],
        "Alert volume in a typical mid-size fleet is not overwhelming on its own. "
        "The failure mode is queue latency - an alert waits for an available "
        "analyst, not for triage to finish - and inconsistent documentation, "
        "with most detections closed leaving no written record at all. A closed "
        "detection with no report leaves nothing to query later: no way to see "
        "repeat offenders, coverage gaps, or trend shifts.",
    )

    # 3 - before / after structural comparison
    deck.flow_compare_slide(
        "What actually changes",
        "Not “faster” - a different set of steps, with different owners",
        "Manual path",
        [
            "Detection fires → lands in queue",
            "Sits until an analyst is free (queue latency)",
            "Analyst manually pulls host / user / process context",
            "Analyst researches indicators, applies judgment",
            "Analyst writes the report by hand (or skips it)",
            "Analyst posts to Slack / closes the console",
        ],
        "Automated path",
        [
            "Detection fires → auto-enriched immediately",
            "Claude applies the same triage skill and policy lenses a human uses",
            "Structured verdict report generated (locked schema)",
            "Report posts to Slack automatically",
            "Analyst reviews the verdict (accept, correct, escalate)",
            "Analyst decision is the only manual step left",
        ],
        "What moved: context-gathering and report-writing shift from analyst to "
        "pipeline. Judgment stays human - the agent proposes, it never closes "
        "or actions a detection on its own.",
    )

    # 4 - architecture
    deck.pipeline_slide(
        "Architecture",
        "Outbound-only, read-only by default, no server exposed to the internet",
        [
            ("CrowdStrike Falcon", "Detection fires (Fusion SOAR event trigger), multi-tenant aware"),
            ("Orchestration layer", "AWS Lambda behind a Function URL; enriches via the Falcon API, groups bursts via S3"),
            ("Claude triage agent", "Loads the triage skill, policy rules, and locked report schema; scoped Anthropic API key"),
            ("Slack + S3 archive", "Structured report posts to the security channel; every report persisted for audit (12-month retention)"),
            ("Analyst review", "Human approves before anything acts; response actions (Phase 3) stay gated behind explicit approval"),
        ],
        "All calls are outbound HTTPS from Falcon's cloud or the Lambda - no inbound ports.",
    )

    # 5 - design principles
    deck.numbered_list_slide(
        "Design principles",
        "Three decisions that shape everything else",
        [
            (
                "Human-in-the-loop, always",
                "The pipeline produces verdicts, not actions. Automated containment on a "
                "false positive can take down a production system - every report lands "
                "in front of an analyst first, and response actions stay behind an "
                "explicit approval step.",
            ),
            (
                "Read-only by default",
                "The Falcon API key holds Alerts-Read and Hosts-Read, nothing else. The "
                "Anthropic key can only create messages; the Slack webhook can only post "
                "to one channel. Any write or remediation call requires a separate, "
                "explicitly reviewed credential.",
            ),
            (
                "Detection content is data, never instructions",
                "Attacker-controllable strings (command lines, filenames, domains) are "
                "treated as untrusted input. The system prompt is locked and the report "
                "schema is fixed, so detection content cannot steer the model into a "
                "different output shape.",
            ),
        ],
        dark=True,
    )

    # 6 - status
    deck.table_slide(
        "Status",
        "Honest phase breakdown - not everything here is finished",
        ["Phase", "Scope", "Status"],
        [
            [
                "1 - Report automation",
                "Falcon-native flow: detection → Claude analysis → structured Slack report. Zero servers.",
                "Built and running",
            ],
            [
                "2 - Unified alert reporting",
                "Lambda orchestration: multi-tenant enrichment, burst grouping, locked schema, S3 archive, recurring digest. Infra as Terraform.",
                "In progress",
            ],
            [
                "3 - Supervised response actions",
                "Host containment / hash blocking, proposed by the agent, executed only after human approval.",
                "Planned",
            ],
        ],
    )

    # 7 - what this demonstrates
    deck.body_slide(
        "Co-pilot for detection engineers, not autopilot",
        "The agent proposes; a human still decides",
        "This pipeline is not about replacing a security analyst - it is about "
        "removing the two tasks that don't need human judgment (waiting in a "
        "queue, and writing up what was found) so the analyst's time goes to "
        "the one task that does (deciding what to do about it).\n\n"
        "Illustrative business case: at a typical alert volume for a fleet this "
        "size, recovering even a fraction of the ~30 minutes of hands-on time "
        "per detection pays for the infrastructure (~$15/month, illustrative) "
        "inside the first day of any given month.",
    )

    # 8 - closing
    deck.closing_slide(
        "Falcon → Claude",
        "Detection triage, automated end to end - judgment stays human",
        [
            "GitHub: github.com/StudioMatan/claude-agentic-workspace",
            "LinkedIn: linkedin.com/in/[placeholder]",
        ],
    )

    default_out = os.path.join(
        os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
        "flows", "falcon-claude-soar", "SOAR-Overview.pptx",
    )
    path = deck.save(out_path or default_out)
    print(f"Saved: {path} ({deck.slide_count} slides)")
    return path


if __name__ == "__main__":
    _cli()
